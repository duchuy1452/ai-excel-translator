from pptx import Presentation
from pptx.shapes.group import GroupShape
from pptx.slide import Slide
from typing import Callable
from io import BytesIO

from translators import GeminiTranslator
from utils import logger


class PowerPointProcessor:
    """
    Processes PowerPoint (.pptx) files, extracts text from various elements,
    translates it using a translator, and creates a new PowerPoint file with
    the translated text.
    """

    def __init__(self, translator: GeminiTranslator):
        self.translator = translator

    def process_powerpoint(
        self, file_data: BytesIO, target_lang: str, progress_callback: Callable = None
    ) -> BytesIO:
        """
        Processes a PowerPoint file, translates its content, and returns the translated file as a BytesIO stream.

        Args:
            file_data (BytesIO): The input PowerPoint file as a BytesIO stream.
            target_lang (str): The target language for translation.
            progress_callback (Callable, optional): A callback function to track progress.

        Returns:
            BytesIO: The translated PowerPoint file as a BytesIO stream.
        """
        presentation = Presentation(file_data)
        total_slides = len(presentation.slides)
        processed_slides = 0

        for slide_idx, slide in enumerate(presentation.slides):
            texts = self._collect_texts_from_slide(slide, target_lang)
            if not texts:
                logger.info("Slide %s has no text. Skipping", slide_idx + 1)
                continue
            logger.info("Process slide %s", slide_idx + 1)
            translated_texts = self.translator.translate_batch(texts, target_lang)
            self._replace_text_in_slide(slide, translated_texts)

            processed_slides += 1
            if progress_callback:
                progress_callback(processed_slides, total_slides)

        output = BytesIO()
        presentation.save(output)
        output.seek(0)
        return output

    def _collect_texts_from_slide(self, slide: Slide, target_lang) -> list[str]:
        """Collects texts from a slide"""
        texts = []
        for shape in slide.shapes:
            texts.extend(self._extract_text_from_shape(shape, target_lang))

        return texts

    def _extract_text_from_shape(self, shape, target_lang) -> list[str]:
        """
        Extracts text from a shape in a PowerPoint slide.

        Handles different shape types including placeholders, text boxes, tables, and groups.

        Args:
            shape: The shape to extract text from.

        Returns:
            A list of strings containing the extracted text.
        """
        texts = []

        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    texts.append(run.text)
        elif isinstance(shape, GroupShape):
            for s in shape.shapes:
                texts.extend(self._extract_text_from_shape(s, target_lang))
        elif shape.has_table:
            self._translate_table(shape, target_lang)
        return texts

    def _translate_table(self, shape, target_lang):
        table = shape.table  # Get the table
        table_cells = []
        table_texts = []
        for row in table.rows:
            for cell in row.cells:
                if cell.text:
                    table_texts.append(cell.text)
                    table_cells.append(cell)

        translated_texts = self.translator.translate_batch(table_texts, target_lang)
        for index, translated_text in enumerate(translated_texts):
            table_cells[index].text = translated_text

    def _replace_text_in_slide(self, slide: Slide, translated_texts: list[str]):
        """Replaces texts in a slide"""
        text_index = 0
        for shape in slide.shapes:
            text_index = self._replace_text_in_shape(
                shape, translated_texts, text_index
            )

    def _replace_text_in_shape(
        self, shape, translated_texts: list[str], text_index: int
    ) -> int:
        """
        Replace text in a shape with translated text.

        Args:
            shape: The shape to replace text in.
            translated_texts: A list of translated strings.
            text_index: The current index in the translated_texts list.

        Returns:
            The updated text_index after replacing text in the shape.
        """
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if text_index < len(translated_texts):
                        run.text = translated_texts[text_index]
                        text_index += 1
        elif isinstance(shape, GroupShape):
            for s in shape.shapes:
                text_index = self._replace_text_in_shape(
                    s, translated_texts, text_index
                )
        return text_index
